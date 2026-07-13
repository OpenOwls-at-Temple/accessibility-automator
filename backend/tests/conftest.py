"""Backend test fixtures: an in-memory user DB + a TestClient backed by temp
per-test storage, plus helpers to register/log in, build a tiny PPTX, and poll
a remediation job.

Auth is DB-backed + JWT bearer: ``login`` registers a user (the allowlist),
dev-logs-in, and attaches the ``Authorization`` header to the client so the
existing request call sites need no per-call headers.
"""

from __future__ import annotations

import io
import time

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

import backend.app.models  # noqa: F401  (register models on Base)
from backend.app.core.database import Base, get_db
from backend.app.main import create_app
from backend.app.models.user import User
from backend.app.settings import Settings
from remediator.config import load_config

API = "/api/v1"


@pytest.fixture
def db_session():
    """Fresh in-memory SQLite user database per test."""
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(engine)
    Session = sessionmaker(bind=engine, autoflush=False, expire_on_commit=False)
    session = Session()
    try:
        yield session
    finally:
        session.close()
        engine.dispose()


@pytest.fixture
def app(tmp_path, db_session):
    settings = Settings(
        google_client_id="test-client-id",
        jwt_secret="test-secret",
        storage_dir=tmp_path / "storage",
        cors_origins=["http://localhost:5173"],
        environment="local",
    )
    application = create_app(settings=settings, config=load_config())

    def override_get_db():
        yield db_session

    application.dependency_overrides[get_db] = override_get_db
    return application


@pytest.fixture
def client(app):
    return TestClient(app)


@pytest.fixture
def pptx_bytes() -> bytes:
    """A deck with no document title and an untitled slide (fails P1, P2)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[1].text = "Mitochondria are the powerhouse of the cell."
    prs.core_properties.title = ""
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def register(
    db_session,
    email: str = "prof@temple.edu",
    *,
    is_admin: bool = False,
    name: str = "Prof",
    is_active: bool = True,
) -> User:
    """Add a user to the allowlist (what an admin invite does)."""
    user = User(email=email.strip().lower(), name=name, is_admin=is_admin, is_active=is_active)
    db_session.add(user)
    db_session.commit()
    db_session.refresh(user)
    return user


def login(
    client: TestClient, db_session, email: str = "prof@temple.edu", *, is_admin: bool = False
):
    """Register + dev-login + attach the bearer header to this client."""
    register(db_session, email, is_admin=is_admin)
    resp = client.post(f"{API}/auth/dev-login", json={"email": email})
    assert resp.status_code == 200, resp.text
    client.headers["Authorization"] = f"Bearer {resp.json()['access_token']}"
    return resp.json()


def upload(client: TestClient, group: str, filename: str, data: bytes):
    return client.post(
        f"{API}/groups/{group}/files",
        files={"files": (filename, data, "application/octet-stream")},
    )


def wait_for_job(client: TestClient, job_id: str, timeout: float = 20.0) -> dict:
    deadline = time.time() + timeout
    while time.time() < deadline:
        data = client.get(f"{API}/jobs/{job_id}").json()
        if data["status"] in ("done", "error"):
            return data
        time.sleep(0.1)
    raise AssertionError(f"Job {job_id} did not finish in {timeout}s")
