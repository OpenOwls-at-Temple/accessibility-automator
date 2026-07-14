from pptx import Presentation
from pptx.util import Inches

from remediator.config import load_config
from remediator.handlers.pptx_handler import PptxHandler, get_alt_text
from remediator.models import ACTION_AI_FIXED, ACTION_NOT_FIXED
from remediator.pipeline import remediate_file
from remediator.rules.pptx_rules import audit_pptx


def _check_passed(results, check_id):
    matching = [r for r in results if r.check_id == check_id]
    return bool(matching) and all(r.passed for r in matching)


def test_pipeline_with_ai_reaches_100_and_no_placeholders(bad_pptx, tmp_path, fake_provider):
    out = tmp_path / "bad_a11y.pptx"
    report = remediate_file(bad_pptx, out, cfg=load_config(), provider=fake_provider)

    assert report.pre_fix_score < report.post_fix_score
    assert report.post_fix_score == 100
    assert report.truly_remediated_score == 100  # all genuinely fixed
    assert report.placeholder_fixes == []

    # Re-audit the written output independently.
    post = audit_pptx(PptxHandler(out))
    for cid in ("P1", "P2", "P3", "P5", "P6"):
        assert _check_passed(post, cid)


def test_pipeline_without_llm_uses_placeholders(bad_pptx, tmp_path):
    out = tmp_path / "bad_a11y.pptx"
    report = remediate_file(bad_pptx, out, cfg=load_config(), provider=None)

    # Checker passes, but truly-remediated is lower because P2/P3 are placeholders.
    assert report.post_fix_score == 100
    assert report.truly_remediated_score < 100
    placeholder_checks = {f.check_id for f in report.placeholder_fixes}
    assert {"P2", "P3"} <= placeholder_checks


def test_original_input_is_never_modified(bad_pptx, tmp_path, fake_provider):
    out = tmp_path / "bad_a11y.pptx"
    remediate_file(bad_pptx, out, cfg=load_config(), provider=fake_provider)

    # The input must still fail its original checks.
    original = audit_pptx(PptxHandler(bad_pptx))
    assert not _check_passed(original, "P3")
    assert not _check_passed(original, "P1")


def test_ai_alt_text_is_written(bad_pptx, tmp_path, fake_provider):
    out = tmp_path / "bad_a11y.pptx"
    report = remediate_file(bad_pptx, out, cfg=load_config(), provider=fake_provider)

    alt_fixes = [f for f in report.fixes if f.check_id == "P3"]
    assert alt_fixes and alt_fixes[0].action == ACTION_AI_FIXED

    # The alt text is actually present in the saved file.
    from remediator.rules.pptx_rules import iter_pictures

    pics = list(iter_pictures(PptxHandler(out).presentation))
    assert get_alt_text(pics[0][1]) == "A blue square test image"


def test_slide_without_title_placeholder_does_not_crash(tmp_path, fake_provider):
    """A slide on the Blank layout has no title placeholder to clone from.

    The title fixer must degrade gracefully (P2 left unfixed) rather than raise
    — a remediation job must always finish and produce a valid output file.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank: no title placeholder
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb.text_frame.text = "Content with no title placeholder anywhere"
    src = tmp_path / "notitle.pptx"
    prs.save(str(src))

    out = tmp_path / "notitle_a11y.pptx"
    report = remediate_file(src, out, cfg=load_config(), provider=fake_provider)

    # The job completed and wrote a valid, reopenable deck.
    assert out.exists()
    assert len(Presentation(str(out)).slides._sldIdLst) == 1
    # P2 could not be placed (no placeholder to hold a title) → honestly unfixed.
    p2 = [f for f in report.fixes if f.check_id == "P2"]
    assert p2 and p2[0].action == ACTION_NOT_FIXED and not p2[0].success


def test_title_with_soft_line_break_does_not_crash(tmp_path):
    """A slide title with a soft line break (U+000B) must not crash the fixer.

    PowerPoint stores in-title line breaks as a vertical tab, which is legal in
    a run but illegal to write into ``core_properties.title`` — lxml raises. The
    derived document title must be sanitized so the job finishes and the output
    is a valid deck (regression: a real 79-slide lecture crashed here).
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only: has a title
    slide.shapes.title.text = "Recursively Defined Sets and \x0bStructural Induction"
    src = tmp_path / "breaky.pptx"
    prs.save(str(src))

    out = tmp_path / "breaky_a11y.pptx"
    report = remediate_file(src, out, cfg=load_config(), provider=None)

    # Completed and wrote a valid, reopenable deck with a clean, control-char-free title.
    assert out.exists()
    written = Presentation(str(out)).core_properties.title
    assert written == "Recursively Defined Sets and Structural Induction"
    assert not any(ord(c) < 0x20 and c not in "\t\n\r" for c in written)
    p1 = [f for f in report.fixes if f.check_id == "P1"]
    assert p1 and p1[0].success


def test_decorative_image_is_marked(bad_pptx, tmp_path, decorative_provider):
    out = tmp_path / "bad_a11y.pptx"
    report = remediate_file(bad_pptx, out, cfg=load_config(), provider=decorative_provider)

    post = audit_pptx(PptxHandler(out))
    assert _check_passed(post, "P3")  # decorative images satisfy the alt-text check
    p3 = [f for f in report.fixes if f.check_id == "P3"][0]
    assert p3.action == ACTION_AI_FIXED and "decorative" in p3.detail.lower()
