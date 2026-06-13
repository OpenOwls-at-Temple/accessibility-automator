"""Shared test fixtures: programmatically built known-bad / known-good PPTX
files (so no binaries live in git) plus a fake LLM provider.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from remediator.handlers.pptx_handler import set_alt_text
from remediator.llm.provider import CaptionResult, TitleResult


def _png_bytes() -> io.BytesIO:
    img = Image.new("RGB", (64, 64), (120, 160, 200))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def _set_lang(run, lang="en-US") -> None:
    run._r.get_or_add_rPr().set("lang", lang)


@pytest.fixture
def png_bytes() -> bytes:
    return _png_bytes().getvalue()


@pytest.fixture
def bad_pptx(tmp_path):
    """A deck that fails P1, P2, P3, P5, P6."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    # P2: leave the title placeholder empty, but give the slide body text.
    slide.placeholders[1].text = "Light reactions occur in the thylakoid membrane."
    # P3: image with no alt text.
    slide.shapes.add_picture(_png_bytes(), Inches(1), Inches(3), Inches(2), Inches(2))
    # P5: table with no header row.
    table = slide.shapes.add_table(2, 2, Inches(4), Inches(3), Inches(4), Inches(1)).table
    table.first_row = False
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    # P1: empty document title. P6: runs carry no language attribute by default.
    prs.core_properties.title = ""
    path = tmp_path / "bad.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def good_pptx(tmp_path):
    """A fully accessible deck — should score 100."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Photosynthesis Overview"
    body = slide.placeholders[1]
    body.text = "Light reactions occur in the thylakoid membrane."
    for para in body.text_frame.paragraphs:
        for run in para.runs:
            _set_lang(run)
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            _set_lang(run)
    pic = slide.shapes.add_picture(_png_bytes(), Inches(1), Inches(3), Inches(2), Inches(2))
    set_alt_text(pic, "Diagram of the light-dependent reactions")
    table = slide.shapes.add_table(2, 2, Inches(4), Inches(3), Inches(4), Inches(1)).table
    table.first_row = True
    prs.core_properties.title = "Biology Lecture 1"
    path = tmp_path / "good.pptx"
    prs.save(str(path))
    return path


class FakeProvider:
    """High-confidence stand-in for the LLM (no network)."""

    def __init__(self, decorative=False):
        self.decorative = decorative

    def caption_image(self, image_bytes, mime_type="image/png", context=""):
        if self.decorative:
            return CaptionResult(alt_text="", decorative=True, confidence=0.9)
        return CaptionResult(alt_text="A blue square test image", decorative=False, confidence=0.9)

    def suggest_title(self, slide_text):
        return TitleResult(title="Light Reactions", confidence=0.9)


@pytest.fixture
def fake_provider():
    return FakeProvider()


@pytest.fixture
def decorative_provider():
    return FakeProvider(decorative=True)


# ── PDF fixtures (built with reportlab; encrypted via pikepdf) ──


def _build_text_pdf(path) -> None:
    """An untagged PDF with text but no title and no language (fails D1/D2/D12)."""
    import pikepdf
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    tmp = str(path) + ".tmp.pdf"
    c = canvas.Canvas(tmp, pagesize=letter)
    c.drawString(72, 720, "Photosynthesis converts light into chemical energy.")
    c.drawString(72, 700, "Light reactions occur in the thylakoid membrane.")
    c.showPage()
    c.save()
    # reportlab writes a default /Title of "untitled"; strip it so D2 truly fails.
    with pikepdf.open(tmp) as pdf:
        if "/Title" in pdf.docinfo:
            del pdf.docinfo["/Title"]
        pdf.save(str(path))
    Path(tmp).unlink()


def _build_image_only_pdf(path, image_path) -> None:
    """A page with only an image and no text (a 'scanned' page: fails D8/D3)."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawImage(str(image_path), 72, 400, width=200, height=200)
    c.showPage()
    c.save()


@pytest.fixture
def bad_pdf(tmp_path):
    path = tmp_path / "bad.pdf"
    _build_text_pdf(path)
    return path


@pytest.fixture
def image_pdf(tmp_path):
    img = tmp_path / "pic.png"
    Image.new("RGB", (200, 200), (90, 140, 200)).save(str(img))
    path = tmp_path / "scanned.pdf"
    _build_image_only_pdf(path, img)
    return path


@pytest.fixture
def encrypted_pdf(tmp_path):
    import pikepdf

    src = tmp_path / "plain.pdf"
    _build_text_pdf(src)
    path = tmp_path / "encrypted.pdf"
    with pikepdf.open(str(src)) as pdf:
        pdf.save(str(path), encryption=pikepdf.Encryption(owner="o", user="u"))
    return path
