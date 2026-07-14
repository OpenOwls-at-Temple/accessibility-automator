"""The remediation pipeline: audit -> fix -> re-score -> report.

Format-agnostic orchestration. A per-format entry in ``_REGISTRY`` binds a
handler, an audit function, and a fix function; adding PDF (or Word in Phase 2)
means adding one registry entry, not changing this file.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from remediator.config import Config, load_config
from remediator.fixers import pdf_fixer, pptx_fixer
from remediator.handlers.pdf_handler import PdfHandler
from remediator.handlers.pptx_handler import PptxHandler
from remediator.llm.provider import build_provider
from remediator.models import ACTION_PLACEHOLDER, FileReport
from remediator.rules import pdf_rules, pptx_rules
from remediator.scorer import compute_score


@dataclass
class SuggestionItem:
    """One AI-generated suggestion for human review before it is written."""

    check_id: str  # e.g. "P3"
    element_ref: str  # e.g. "Slide 4 / Picture 7"
    suggestion_type: str  # "alt_text" | "slide_title"
    draft_text: str  # what the AI (or placeholder) would write
    is_placeholder: bool  # True when we had no real LLM to call


# extension -> (handler class, audit fn, fix fn)
_REGISTRY = {
    ".pptx": (PptxHandler, pptx_rules.audit_pptx, pptx_fixer.fix_pptx),
    ".pdf": (PdfHandler, pdf_rules.audit_pdf, pdf_fixer.fix_pdf),
}


def supported_extensions() -> tuple[str, ...]:
    return tuple(_REGISTRY)


def scan_file(
    input_path: str | Path,
    cfg: Config | None = None,
    provider="__build__",
) -> list[SuggestionItem]:
    """Audit a file and generate AI draft suggestions WITHOUT writing anything.

    Returns one SuggestionItem per element that needs a human decision (missing
    alt text, missing slide title). The caller shows these to the user for
    editing before calling remediate_file with the approved overrides.
    """
    input_path = Path(input_path)
    cfg = cfg or load_config()
    ext = input_path.suffix.lower()
    if ext not in _REGISTRY:
        raise NotImplementedError(f"Unsupported file type {ext!r}")
    handler_cls, _audit_fn, _ = _REGISTRY[ext]

    if provider == "__build__":
        provider = build_provider(cfg.llm)

    handler = handler_cls(input_path)

    suggestions: list[SuggestionItem] = []

    if ext == ".pptx":
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from remediator.handlers.pptx_handler import (
            get_alt_text,
            is_decorative,
            is_meaningful_alt_text,
        )
        from remediator.llm.provider import LLMError
        from remediator.rules import pptx_rules as _pr

        prs = handler.presentation
        for s_idx, slide in enumerate(prs.slides, start=1):
            # P2 — missing slide title
            title_shape = slide.shapes.title
            if title_shape is None or not (title_shape.text or "").strip():
                slide_text = " ".join(
                    s.text.strip()
                    for s in _pr._walk_shapes(slide.shapes)
                    if s.has_text_frame and s.text.strip()
                )
                draft = f"[Slide {s_idx} Title]"
                is_ph = True
                if provider and cfg.llm.suggest_titles and slide_text.strip():
                    try:
                        result = provider.suggest_title(slide_text)
                        if result.title and result.confidence >= cfg.llm.confidence_threshold:
                            draft = result.title
                            is_ph = False
                    except (LLMError, Exception):
                        pass
                suggestions.append(
                    SuggestionItem(
                        check_id="P2",
                        element_ref=_pr._slide_label(s_idx),
                        suggestion_type="slide_title",
                        draft_text=draft,
                        is_placeholder=is_ph,
                    )
                )

            # P3 — images missing alt text
            for shape in _pr._walk_shapes(slide.shapes):
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                if is_decorative(shape) or is_meaningful_alt_text(get_alt_text(shape)):
                    continue
                context = slide_text if "slide_text" in dir() else ""
                draft = cfg.signoff.placeholder_alt_text
                is_ph = True
                if provider and cfg.llm.generate_alt_text:
                    try:
                        image = shape.image
                        result = provider.caption_image(image.blob, image.content_type, context)
                        if (
                            not result.decorative
                            and result.alt_text
                            and result.confidence >= cfg.llm.confidence_threshold
                        ):
                            draft = result.alt_text
                            is_ph = False
                    except (LLMError, Exception):
                        pass
                suggestions.append(
                    SuggestionItem(
                        check_id="P3",
                        element_ref=_pr._ref(s_idx, shape),
                        suggestion_type="alt_text",
                        draft_text=draft,
                        is_placeholder=is_ph,
                    )
                )

    return suggestions


def remediate_file(
    input_path: str | Path,
    output_path: str | Path,
    cfg: Config | None = None,
    provider="__build__",
    overrides: dict[str, str] | None = None,
) -> FileReport:
    """Remediate one file and return its before/after report.

    ``provider`` defaults to building one from config + environment; pass an
    explicit provider (or ``None`` to force placeholder-only) for tests.
    """
    input_path, output_path = Path(input_path), Path(output_path)
    cfg = cfg or load_config()
    ext = input_path.suffix.lower()
    if ext not in _REGISTRY:
        raise NotImplementedError(
            f"Unsupported file type {ext!r}; supported: {supported_extensions()}"
        )
    handler_cls, audit_fn, fix_fn = _REGISTRY[ext]

    if provider == "__build__":
        provider = build_provider(cfg.llm)

    # 1. Audit the input -> pre-fix score.
    handler = handler_cls(input_path)
    pre_audit = audit_fn(handler)
    pre_score = compute_score(pre_audit, cfg.scoring)

    # 2. Fix, then write the remediated copy (originals are never touched).
    fixes = fix_fn(handler, pre_audit, cfg, provider, overrides or {})
    output_path.parent.mkdir(parents=True, exist_ok=True)
    handler.save(output_path)

    # 3. Re-audit the output -> checker-passing + truly-remediated scores.
    post_handler = handler_cls(output_path)
    post_audit = audit_fn(post_handler)
    post_score = compute_score(post_audit, cfg.scoring)
    placeholder_checks = {f.check_id for f in fixes if f.action == ACTION_PLACEHOLDER}
    truly_score = compute_score(post_audit, cfg.scoring, placeholder_checks)

    return FileReport(
        file_name=input_path.name,
        file_type=handler.file_type,
        pre_fix_audit=pre_audit,
        fixes=fixes,
        post_fix_audit=post_audit,
        pre_fix=pre_score,
        post_fix=post_score,
        truly_remediated=truly_score,
    )
