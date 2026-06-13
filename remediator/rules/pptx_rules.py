"""PPTX audit checks P1–P13 (approximating the YuJa Panorama checklist).

Each check emits zero or more :class:`AuditResult`. Element-scoped checks
(alt text, tables, fonts) emit one result per relevant element and emit nothing
when no such element exists — the check is simply absent from the score when it
does not apply. Document-scoped checks (title, language) always emit.

See ai_specs/architecture-planning.md §Panorama rule sets for severities and
strategies. The fixer in ``fixers/pptx_fixer.py`` mirrors these check ids.
"""

from __future__ import annotations

from collections.abc import Iterator

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from remediator.handlers.pptx_handler import (
    PptxHandler,
    get_alt_text,
    is_decorative,
    is_meaningful_alt_text,
)
from remediator.models import (
    SEVERITY_DISABLED,
    SEVERITY_MAJOR,
    SEVERITY_MINOR,
    SEVERITY_SEVERE,
    AuditResult,
)

_GENERIC_LINK_TEXT = {"click here", "here", "link", "read more", "more", "this", "click"}


# ── shape / run traversal helpers ──


def _walk_shapes(shapes) -> Iterator:
    """Yield every shape, descending into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _slide_label(idx: int) -> str:
    return f"Slide {idx}"


def _ref(idx: int, shape) -> str:
    return f"{_slide_label(idx)} / {shape.name}"


def iter_pictures(presentation):
    """Yield (slide_index, shape) for every picture in the deck."""
    for s_idx, slide in enumerate(presentation.slides, start=1):
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield s_idx, shape


def iter_runs(presentation):
    """Yield (slide_index, shape, run) for every text run in the deck."""
    for s_idx, slide in enumerate(presentation.slides, start=1):
        for shape in _walk_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    yield s_idx, shape, run


def _run_lang(run) -> str | None:
    rpr = run._r.find(qn("a:rPr"))
    return rpr.get("lang") if rpr is not None else None


def _relative_luminance(rgb) -> float:
    def channel(c: int) -> float:
        cs = c / 255.0
        return cs / 12.92 if cs <= 0.03928 else ((cs + 0.055) / 1.055) ** 2.4

    r, g, b = rgb[0], rgb[1], rgb[2]
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def _contrast_ratio(rgb1, rgb2) -> float:
    l1, l2 = _relative_luminance(rgb1), _relative_luminance(rgb2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _slide_background_rgb(slide):
    """Return the slide's solid background colour, or None if not resolvable."""
    try:
        fill = slide.background.fill
        rgb = fill.fore_color.rgb  # raises unless a solid RGB fill
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


# ── individual checks ──


def check_p1_document_title(presentation) -> list[AuditResult]:
    title = (presentation.core_properties.title or "").strip()
    return [
        AuditResult(
            "P1",
            SEVERITY_MAJOR,
            passed=bool(title),
            element_ref="Document",
            detail="Document title set" if title else "Document has no title metadata",
            recommendation="" if title else "Set the presentation title in file properties.",
        )
    ]


def check_p2_slide_titles(presentation) -> list[AuditResult]:
    results = []
    for s_idx, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        text = (title_shape.text or "").strip() if title_shape is not None else ""
        results.append(
            AuditResult(
                "P2",
                SEVERITY_MAJOR,
                passed=bool(text),
                element_ref=_slide_label(s_idx),
                detail="Slide has a title" if text else "Slide has no title",
                recommendation="" if text else "Add a descriptive slide title.",
            )
        )
    return results


def check_p3_image_alt_text(presentation) -> list[AuditResult]:
    results = []
    for s_idx, shape in iter_pictures(presentation):
        if is_decorative(shape):
            passed, detail = True, "Marked decorative"
        else:
            alt = get_alt_text(shape)
            passed = is_meaningful_alt_text(alt)
            if passed:
                detail = f"Alt text: {alt!r}"
            elif alt:
                detail = f"Alt text not meaningful (filename): {alt!r}"
            else:
                detail = "Image has no alt text"
        results.append(
            AuditResult(
                "P3",
                SEVERITY_MAJOR,
                passed=passed,
                element_ref=_ref(s_idx, shape),
                detail=detail,
                recommendation="" if passed else "Add alt text or mark the image decorative.",
            )
        )
    return results


def check_p4_contrast(presentation) -> list[AuditResult]:
    """Report-only, conservative: only flags when both text and background
    colours are explicitly resolvable, to avoid false positives."""
    results = []
    for s_idx, slide in enumerate(presentation.slides, start=1):
        bg = _slide_background_rgb(slide)
        if bg is None:
            continue
        for shape in _walk_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb = run.font.color.rgb
                        fg = (rgb[0], rgb[1], rgb[2])
                    except Exception:
                        continue
                    size = run.font.size
                    large = size is not None and size >= Pt(18)
                    threshold = 3.0 if large else 4.5
                    ratio = _contrast_ratio(fg, bg)
                    if ratio < threshold:
                        results.append(
                            AuditResult(
                                "P4",
                                SEVERITY_MAJOR,
                                passed=False,
                                element_ref=_ref(s_idx, shape),
                                detail=f"Contrast {ratio:.2f}:1 below {threshold}:1",
                                recommendation="Increase text/background contrast.",
                            )
                        )
    return results


def check_p5_table_headers(presentation) -> list[AuditResult]:
    results = []
    for s_idx, slide in enumerate(presentation.slides, start=1):
        for shape in _walk_shapes(slide.shapes):
            if not shape.has_table:
                continue
            passed = bool(shape.table.first_row)
            results.append(
                AuditResult(
                    "P5",
                    SEVERITY_MAJOR,
                    passed=passed,
                    element_ref=_ref(s_idx, shape),
                    detail="Header row set" if passed else "Table has no header row",
                    recommendation="" if passed else "Mark the first row as a header.",
                )
            )
    return results


def check_p6_language(presentation) -> list[AuditResult]:
    runs = list(iter_runs(presentation))
    if not runs:
        return []
    has_lang = any(_run_lang(run) for _, _, run in runs)
    return [
        AuditResult(
            "P6",
            SEVERITY_MINOR,
            passed=has_lang,
            element_ref="Document",
            detail="Language specified" if has_lang else "No language specified",
            recommendation="" if has_lang else "Set the document language.",
        )
    ]


def check_p7_incorrect_language(presentation) -> list[AuditResult]:
    """Flag only clearly malformed/empty language tags to avoid false positives
    (a slide may be legitimately authored in another language)."""
    results = []
    for s_idx, shape, run in iter_runs(presentation):
        lang = _run_lang(run)
        if lang is not None and not lang.strip():
            results.append(
                AuditResult(
                    "P7",
                    SEVERITY_MINOR,
                    passed=False,
                    element_ref=_ref(s_idx, shape),
                    detail="Empty language tag",
                    recommendation="Set a valid language tag.",
                )
            )
    return results


def check_p8_link_text(presentation) -> list[AuditResult]:
    results = []
    for s_idx, shape, run in iter_runs(presentation):
        addr = run.hyperlink.address
        if not addr:
            continue
        text = (run.text or "").strip()
        bare = text == addr or text.lower().startswith(("http://", "https://", "www."))
        generic = text.lower() in _GENERIC_LINK_TEXT
        if bare or generic or not text:
            results.append(
                AuditResult(
                    "P8",
                    SEVERITY_MINOR,
                    passed=False,
                    element_ref=_ref(s_idx, shape),
                    detail=f"Non-descriptive link text: {text!r} -> {addr}",
                    recommendation="Replace with descriptive link text.",
                )
            )
    return results


def check_p9_font_size(presentation) -> list[AuditResult]:
    results = []
    for s_idx, shape, run in iter_runs(presentation):
        size = run.font.size
        if size is not None and size < Pt(9) and (run.text or "").strip():
            results.append(
                AuditResult(
                    "P9",
                    SEVERITY_MINOR,
                    passed=False,
                    element_ref=_ref(s_idx, shape),
                    detail=f"Font size {size.pt:.1f}pt below 9pt",
                    recommendation="Increase font size to at least 9pt.",
                )
            )
    return results


def check_p10_reading_order(presentation) -> list[AuditResult]:
    """Flag slides whose shape order does not match top-to-bottom,
    left-to-right position (a proxy for reading/tab order)."""
    results = []
    for s_idx, slide in enumerate(presentation.slides, start=1):
        positioned = [s for s in slide.shapes if s.top is not None and s.left is not None]
        if len(positioned) < 2:
            continue
        current = [id(s._element) for s in positioned]
        ordered = [id(s._element) for s in sorted(positioned, key=lambda s: (s.top, s.left))]
        passed = current == ordered
        results.append(
            AuditResult(
                "P10",
                SEVERITY_MINOR,
                passed=passed,
                element_ref=_slide_label(s_idx),
                detail=(
                    "Reading order matches layout"
                    if passed
                    else "Shape order does not follow visual layout"
                ),
                recommendation="" if passed else "Reset reading order top-to-bottom.",
            )
        )
    return results


def check_p11_p12_integrity(presentation) -> list[AuditResult]:
    """If we opened the file, it is neither malformed nor encrypted."""
    return [
        AuditResult("P11", SEVERITY_SEVERE, True, "Document", "Document is well-formed"),
        AuditResult("P12", SEVERITY_SEVERE, True, "Document", "Document is not encrypted"),
    ]


def check_p13_format(handler: PptxHandler) -> list[AuditResult]:
    outdated = handler.path.suffix.lower() == ".ppt"
    return [
        AuditResult(
            "P13",
            SEVERITY_DISABLED,
            passed=not outdated,
            element_ref="Document",
            detail="Modern .pptx format" if not outdated else "Outdated .ppt format",
            recommendation="" if not outdated else "Convert to .pptx.",
        )
    ]


def audit_pptx(handler: PptxHandler) -> list[AuditResult]:
    """Run every PPTX check and return the combined audit."""
    prs = handler.presentation
    results: list[AuditResult] = []
    results += check_p1_document_title(prs)
    results += check_p2_slide_titles(prs)
    results += check_p3_image_alt_text(prs)
    results += check_p4_contrast(prs)
    results += check_p5_table_headers(prs)
    results += check_p6_language(prs)
    results += check_p7_incorrect_language(prs)
    results += check_p8_link_text(prs)
    results += check_p9_font_size(prs)
    results += check_p10_reading_order(prs)
    results += check_p11_p12_integrity(prs)
    results += check_p13_format(handler)
    return results
